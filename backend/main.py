import os
import tempfile
import uuid

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from dotenv import load_dotenv
from pptx import Presentation
from groq import AsyncGroq
import azure.cognitiveservices.speech as speechsdk

from sqlalchemy import create_engine, Column, Integer, Text, DateTime, String
from sqlalchemy.orm import declarative_base, sessionmaker
from datetime import datetime, timezone

load_dotenv()

app = FastAPI()

os.makedirs(os.path.join("static", "slides"), exist_ok=True)
app.mount("/static", StaticFiles(directory="static"), name="static")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Database Setup
SQLALCHEMY_DATABASE_URL = "sqlite:///./tts_history.db"
engine = create_engine(SQLALCHEMY_DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

class SynthesisHistory(Base):
    __tablename__ = "synthesis_history"
    id = Column(Integer, primary_key=True, index=True)
    ssml = Column(Text, nullable=False)
    created_at = Column(DateTime, default=lambda: datetime.now(timezone.utc))

class SharedProject(Base):
    __tablename__ = "shared_projects"
    id = Column(String, primary_key=True, index=True)
    project_data = Column(Text, nullable=False)
    created_at = Column(DateTime, default=lambda: datetime.now(timezone.utc))

Base.metadata.create_all(bind=engine)

GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")
AZURE_SPEECH_KEY = os.getenv("AZURE_SPEECH_KEY", "")
AZURE_SPEECH_ENDPOINT = os.getenv("AZURE_SPEECH_ENDPOINT", "https://saisinare19-8483-resource.cognitiveservices.azure.com/")

if GROQ_API_KEY:
    groq_client = AsyncGroq(api_key=GROQ_API_KEY)
    has_groq = True
else:
    has_groq = False

class ScriptRequest(BaseModel):
    text: str
    previous_context: str = ""

class AudioRequest(BaseModel):
    ssml: str



@app.post("/api/upload-ppt")
async def upload_ppt(file: UploadFile = File(...)):
    if not file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")
    
    try:
        content = await file.read()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        prs = Presentation(tmp_path)
        slides_data = []

        session_id = str(uuid.uuid4())
        abs_output_dir = os.path.abspath(os.path.join("static", "slides", session_id))
        os.makedirs(abs_output_dir, exist_ok=True)
        
        try:
            import aspose.slides as slides
            import aspose.pydrawing as drawing
            
            abs_pptx_path = os.path.abspath(tmp_path)
            with slides.Presentation(abs_pptx_path) as presentation:
                for idx, slide in enumerate(presentation.slides):
                    img_name = f"slide_{idx + 1}.png"
                    img_path = os.path.join(abs_output_dir, img_name)
                    
                    bmp = slide.get_image(1.0, 1.0)
                    bmp.save(img_path)
        except Exception as e:
            print(f"Error extracting slide images with aspose: {e}")
            import traceback
            with open(os.path.join(abs_output_dir, "error.txt"), "w") as f:
                f.write(str(e) + "\n" + traceback.format_exc())

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text

            slides_data.append({
                "slide_number": i + 1,
                "text": "\n".join(slide_text),
                "notes": notes,
                "image_url": f"http://localhost:8000/static/slides/{session_id}/slide_{i+1}.png"
            })
            
        os.remove(tmp_path)
        return {"slides": slides_data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate-script")
async def generate_script(req: ScriptRequest):
    if not has_groq:
        raise HTTPException(status_code=500, detail="GROQ_API_KEY is not set in backend/.env")
        
    system_prompt = """
You are an expert presentation script writer. Turn the provided presentation slide text into a highly engaging, conversational script to be read by an AI voice.
You MUST output the script in valid SSML format compatible with Azure Cognitive Services Omni neural voices.

Rules for Azure SSML:
1. Wrap everything in `<speak version="1.0" xmlns="http://www.w3.org/2001/10/synthesis" xmlns:mstts="https://www.w3.org/2001/mstts" xml:lang="en-US">`
2. Next, define the voice: `<voice name="en-US-Ava:DragonHDLatestNeural">`
3. Use the `<mstts:express-as>` tag to add emotions or paralinguistic sounds where appropriate. Supported styles include: 'cheerful', 'sad', 'excited', 'laughing', 'sighing', 'breathing'.
4. Example: <mstts:express-as style="laughing"> That is so funny! </mstts:express-as>
5. Example: <mstts:express-as style="excited"> Let's look at the incredible results! </mstts:express-as>
6. Don't overuse styles, keep it natural. Use them only when it makes sense for a presentation.
7. Only use `<break>` tags when absolutely necessary for a major transition between completely different topics. DO NOT use `<break>` tags between every word, phrase, or normal sentence. When you do use a break, the duration MUST NOT be more than 50ms (e.g., `<break time="50ms"/>`). Keep the pacing natural and energetic.
8. End the SSML correctly with `</voice></speak>`.
9. Do NOT output markdown code blocks (like ```xml). ONLY output the raw valid SSML string.
"""
    try:
        prompt = f"Here is the slide content:\n\n{req.text}\n\n"
        if req.previous_context:
            prompt += f"Context from PREVIOUS slides to maintain presentation flow:\n{req.previous_context[-3000:]}\n\n"
            prompt += "Generate the 30-40 second presentation script in SSML for the CURRENT slide. DO NOT repeat the previous slides' content, just seamlessly continue the presentation as the next part.\n"
        else:
            prompt += "Generate the 30-40 second presentation script in SSML for the CURRENT slide.\n"
            
        completion = await groq_client.chat.completions.create(
            model="openai/gpt-oss-120b",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1024,
        )
        text = completion.choices[0].message.content.strip()
        if text.startswith("```"):
            lines = text.split("\n")
            if len(lines) > 2:
                text = "\n".join(lines[1:-1]).strip()
        return {"ssml": text}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/synthesize")
async def synthesize_audio(req: AudioRequest):
    try:
        speech_config = speechsdk.SpeechConfig(
            subscription=AZURE_SPEECH_KEY, 
            endpoint=AZURE_SPEECH_ENDPOINT
        )
        
        _, temp_path = tempfile.mkstemp(suffix=".wav")
        audio_config = speechsdk.audio.AudioOutputConfig(filename=temp_path)

        speech_synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=audio_config)
        
        result = speech_synthesizer.speak_ssml_async(req.ssml).get()

        if result.reason == speechsdk.ResultReason.SynthesizingAudioCompleted:
            # Log to Database
            db = SessionLocal()
            try:
                new_record = SynthesisHistory(ssml=req.ssml)
                db.add(new_record)
                db.commit()
            except Exception as db_e:
                print(f"Failed to log history: {db_e}")
            finally:
                db.close()
                
            return FileResponse(temp_path, media_type="audio/wav")
        elif result.reason == speechsdk.ResultReason.Canceled:
            cancellation_details = result.cancellation_details
            raise Exception(f"Speech synthesis canceled: {cancellation_details.reason}. Error details: {cancellation_details.error_details}")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/history")
async def get_history():
    db = SessionLocal()
    try:
        records = db.query(SynthesisHistory).order_by(SynthesisHistory.created_at.desc()).all()
        return [
            {
                "id": r.id,
                "ssml": r.ssml,
                "created_at": r.created_at.isoformat()
            }
            for r in records
        ]
    finally:
        db.close()

class ShareRequest(BaseModel):
    project_id: str = None
    project_data: str

@app.post("/api/share")
async def share_project(req: ShareRequest):
    db = SessionLocal()
    try:
        if req.project_id:
            share_id = req.project_id
            existing = db.query(SharedProject).filter(SharedProject.id == share_id).first()
            if existing:
                existing.project_data = req.project_data
            else:
                new_share = SharedProject(id=share_id, project_data=req.project_data)
                db.add(new_share)
        else:
            share_id = uuid.uuid4().hex[:8]
            new_share = SharedProject(id=share_id, project_data=req.project_data)
            db.add(new_share)
        db.commit()
        return {"share_id": share_id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        db.close()

@app.get("/api/share/{share_id}")
async def get_shared_project(share_id: str):
    db = SessionLocal()
    try:
        shared = db.query(SharedProject).filter(SharedProject.id == share_id).first()
        if not shared:
            raise HTTPException(status_code=404, detail="Shared project not found")
        return {"project_data": shared.project_data}
    finally:
        db.close()

import json
@app.get("/api/projects")
async def get_projects():
    db = SessionLocal()
    try:
        projects = db.query(SharedProject).order_by(SharedProject.created_at.desc()).all()
        result = []
        for p in projects:
            try:
                data = json.loads(p.project_data)
                slides = data.get("slides", [])
                image_url = slides[0].get("image_url") if len(slides) > 0 else None
                result.append({
                    "id": p.id,
                    "file_name": data.get("file_name", "Untitled Project"),
                    "slide_count": len(slides),
                    "image_url": image_url,
                    "created_at": p.created_at.isoformat()
                })
            except:
                pass
        return result
    finally:
        db.close()
