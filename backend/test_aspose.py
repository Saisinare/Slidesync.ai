import aspose.slides as slides
import aspose.pydrawing as drawing
import os
from pptx import Presentation

# Create a dummy pptx
prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Hello, World!"
dummy_path = os.path.abspath("dummy2.pptx")
prs.save(dummy_path)

print("Testing Aspose...")
try:
    with slides.Presentation(dummy_path) as presentation:
        for idx, slide in enumerate(presentation.slides):
            img_path = os.path.abspath(f"slide_{idx + 1}_aspose.png")
            print(f"Generating thumbnail for {img_path}...")
            bmp = slide.get_image(1.0, 1.0)
            bmp.save(img_path)
            print(f"Saved {img_path}")
except Exception as e:
    print(f"Aspose error: {e}")
    import traceback
    traceback.print_exc()
