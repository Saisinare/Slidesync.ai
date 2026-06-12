import comtypes.client
import os
from pptx import Presentation

# Create a dummy pptx
prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Hello, World!"
dummy_path = os.path.abspath("dummy.pptx")
prs.save(dummy_path)

# Test comtypes
try:
    print("Starting PowerPoint...")
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    # powerpoint.Visible = False # Sometimes this causes issues
    print("PowerPoint started. Opening presentation...")
    com_prs = powerpoint.Presentations.Open(dummy_path, WithWindow=False)
    print("Presentation opened. Exporting...")
    
    out_dir = os.path.abspath("dummy_out")
    os.makedirs(out_dir, exist_ok=True)
    
    for i in range(1, com_prs.Slides.Count + 1):
        img_path = os.path.join(out_dir, f"slide_{i}.png")
        com_prs.Slides(i).Export(img_path, "PNG")
        print(f"Exported {img_path}")
        
    com_prs.Close()
    powerpoint.Quit()
    print("Success!")
except Exception as e:
    print(f"Error: {e}")
    import traceback
    traceback.print_exc()
